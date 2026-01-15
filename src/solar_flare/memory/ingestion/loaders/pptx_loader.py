"""
PowerPoint document loader using python-pptx.
"""
from pathlib import Path
from typing import BinaryIO, Dict, Any

from .base import BaseDocumentLoader, LoadedDocument, DocumentLoadError


class PptxLoader(BaseDocumentLoader):
    """Load PowerPoint documents using python-pptx."""

    SUPPORTED_EXTENSIONS = [".pptx"]

    def __init__(self, include_notes: bool = True):
        """
        Initialize PowerPoint loader.

        Args:
            include_notes: Whether to extract speaker notes
        """
        self.include_notes = include_notes

    def load(self, file_path: Path) -> LoadedDocument:
        """Load PowerPoint from file path."""
        try:
            from pptx import Presentation
        except ImportError:
            raise DocumentLoadError(
                "python-pptx not installed. Run: pip install python-pptx",
                str(file_path)
            )

        try:
            prs = Presentation(file_path)
            return self._extract_content(prs, str(file_path), file_path.name)
        except Exception as e:
            raise DocumentLoadError(f"Failed to load PowerPoint: {e}", str(file_path), e)

    def load_from_stream(
        self,
        stream: BinaryIO,
        filename: str,
        source_path: str,
    ) -> LoadedDocument:
        """Load PowerPoint from binary stream."""
        try:
            from pptx import Presentation
        except ImportError:
            raise DocumentLoadError(
                "python-pptx not installed. Run: pip install python-pptx",
                source_path
            )

        try:
            prs = Presentation(stream)
            return self._extract_content(prs, source_path, filename)
        except Exception as e:
            raise DocumentLoadError(f"Failed to load PowerPoint stream: {e}", source_path, e)

    def _extract_content(self, prs, source_path: str, filename: str) -> LoadedDocument:
        """Extract content from python-pptx Presentation."""
        text_parts = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"--- Slide {slide_num} ---"]

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)

            # Extract speaker notes if enabled
            if self.include_notes and slide.has_notes_slide:
                try:
                    notes_frame = slide.notes_slide.notes_text_frame
                    if notes_frame and notes_frame.text.strip():
                        slide_text.append(f"[NOTES]: {notes_frame.text}")
                except Exception:
                    pass

            text_parts.append("\n".join(slide_text))

        content = "\n\n".join(text_parts)

        metadata: Dict[str, Any] = {
            "filename": filename,
            "slide_count": len(prs.slides),
        }

        # Try to get title from core properties
        title = self._generate_title(filename)
        try:
            if prs.core_properties and prs.core_properties.title:
                title = prs.core_properties.title
        except Exception:
            pass

        return LoadedDocument(
            content=content,
            title=title,
            source_path=source_path,
            file_format="pptx",
            content_hash=LoadedDocument.compute_hash(content),
            metadata=metadata,
            page_count=len(prs.slides),
        )
